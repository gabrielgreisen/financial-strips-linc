from pptx import Presentation
import pandas as pd
from pptx.dml.color import RGBColor
from helpers.copy_helpers import copy_table_from_template_slide
from helpers.logo_resources import get_logo_file_path_main, get_logo_file_path, get_lincoln_file_path
from helpers.logo_placement import place_logo_on_slide

def financials_layout_one(prs: Presentation, layout_index: int, buyers_chunk_df: pd.DataFrame, start_number: int, brand_api_key):
    
    """
    Adds a slide to the presentation using the specified layout index,
    finds the table, and fills it with financial buyers from the DataFrame slice.

    Parameters
    ----------
    prs : pptx.Presentation
        The loaded Presentation object.

    layout_index : int
        Index of the slide master layout (e.g. 1).

    buyers_chunk_df : pd.DataFrame
        A slice of the DataFrame, typically up to 5 rows.
    """

    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    copy_table_from_template_slide(prs, source_slide_idx=1, target_slide=slide)

    # Find the table shape on the slide
    table = None
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            table_shape = shape
            break
    if table is None:
        raise ValueError(f'No table found on slide layout {layout_index}!')
    
    # Fill table rows
    for i, (_, row) in enumerate(buyers_chunk_df.iterrows()):
        row_idx = i + 1
        # Assumes that columns follow the correct format/order of the columns according to mask; fixed column positions
        type = str(row["primary_type"])
        type2 = str(row["secondary_type"])
        country = str(row["country"])
        br_investment = str(row["brazil_investments"])
        
        def format_number(val):
            try:
                return f"{float(val):,.2f}" # comma as thousand separator + 2 decimals
            except:
                return str(val)

        dry_powder_latam = format_number(row["dry_powder_latam"])
        linc_advised_main = str(row["linc_advised"])

        investment1desc = str(row["investment1_shortdesc"])
        investment2desc = str(row["investment2_shortdesc"])
        investment3desc = str(row["investment3_shortdesc"])

        # Build first column (numbering)
        number = start_number + i
        cell = table.cell(row_idx, 0)
        for para in cell.text_frame.paragraphs:
            if para.runs:
                para.runs[0].text = str(number)
                break
        else:
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = str(number)
    
        # Build second column (potential buyer)
        cell = table.cell(row_idx, 1)
        
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"{country}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"{country}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break
       
            
            
        # Build third column (type)
        cell = table.cell(row_idx, 2)
        is_PE = (type == "PE/Buyout") or ((type2 == "PE/Buyout"))
        is_VC = (type == "Venture Capital") or ((type2 == "Venture Capital"))

        found_index = None
        if is_PE and is_VC:
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"Private Equity\nVenture Capital" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"Private Equity\nVenture Capital"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
        elif is_PE:
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"Private Equity" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"Private Equity"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
        elif is_VC:
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"Venture Capital" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"Venture Capital"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
        else:
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"{type}\n{type2}" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"{type}\n{type2}"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
        

        # Build fourth column (dry powder)
        cell = table.cell(row_idx, 3)
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"{dry_powder_latam}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"{dry_powder_latam}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break
        
        
        # Build fifth column (BR Presence)
        check_mark = "\u2713" # ✓
        cross_mark = "\u2718" # ✘

        if br_investment == "Yes":
            # Figure out how to put the check and the cross formating (its some sort of font)(solved for now)
            cell = table.cell(row_idx, 4)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.text = check_mark
                run.font.color.rgb = RGBColor(0,168,126)
            else:
                cell.text_frame.clear()
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = check_mark
                run.font.color.rgb = RGBColor(0,168,126)
        elif br_investment == "No":
            cell = table.cell(row_idx, 4)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.text = cross_mark
                run.font.color.rgb = RGBColor(192,0,0)
            else:
                cell.text_frame.clear()
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = cross_mark
                run.font.color.rgb = RGBColor(192,0,0)

        # Build sixth column (Investment 1)
        cell = table.cell(row_idx, 5)
        
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"{investment1desc}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"{investment1desc}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break
       
        # Build seventh column (Investment 2)
        cell = table.cell(row_idx, 6)
        
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"{investment2desc}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"{investment2desc}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break

       # Build eigth column (Investment 3)
        cell = table.cell(row_idx, 7)
        
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"{investment3desc}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"{investment3desc}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break
           

        # Add logos to the second column
        logo_file = get_logo_file_path_main(row, brand_api_key=brand_api_key)
        if logo_file:
            place_logo_on_slide(slide, table_shape, table, row_idx, 1, logo_file,
                                width_spacing=0.90, height_spacing=0.60, left_spacing=0.03, top_spacing=0.18)

        # Add logos to the sixth column
        logo_file = get_logo_file_path(row, logo_name_column="investment1_logofile", domain_column="investment1_website", brand_api_key=brand_api_key)
        if logo_file:
            place_logo_on_slide(slide, table_shape, table, row_idx, 5, logo_file,
                                width_spacing=0.90, height_spacing=0.50, left_spacing=0.03, top_spacing=0.18)

        # Add logos to the seventh column
        logo_file = get_logo_file_path(row, logo_name_column="investment2_logofile", domain_column="investment2_website", brand_api_key=brand_api_key)
        if logo_file:
            place_logo_on_slide(slide, table_shape, table, row_idx, 6, logo_file,
                                width_spacing=0.90, height_spacing=0.50, left_spacing=0.03, top_spacing=0.18)

        # Add logos to the eight column
        logo_file = get_logo_file_path(row, logo_name_column="investment3_logofile", domain_column="investment3_website", brand_api_key=brand_api_key)
        if logo_file:
            place_logo_on_slide(slide, table_shape, table, row_idx, 7, logo_file,
                                width_spacing=0.90, height_spacing=0.50, left_spacing=0.03, top_spacing=0.18)

        # Add favicon to investor advised in second column
        if linc_advised_main == "Yes":
            logo_file = get_lincoln_file_path(logo_name="linc_favi", brand_api_key=brand_api_key)
            if logo_file:
                place_logo_on_slide(slide, table_shape, table, row_idx, 1, logo_file,
                                    width_spacing=0.2, height_spacing=0.3, left_spacing=0.90, top_spacing=0.1)