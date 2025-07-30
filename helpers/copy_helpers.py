from pptx import Presentation, slide
from copy import deepcopy

def copy_table_from_template_slide(prs: Presentation, source_slide_idx: int, target_slide: slide.Slide):

    """
    Copies the first table found on a template slide into the target slide.

    Parameters
    ----------
    prs : pptx.Presentation
        The loaded Presentation object that contains all slides.

    source_slide_idx : int
        Index of the slide in prs.slides to use as the table template.
        For example, if your template table is on the first slide, pass 0.

    target_slide : pptx.slide.Slide
        The slide object where the copied table will be inserted.

    Raises
    ------
    ValueError
        If no table is found on the source slide.
    """

    source_slide = prs.slides[source_slide_idx]

    for shape in source_slide.shapes:
        if shape.has_table:
            # Clone the underlying XML element
            new_element = deepcopy(shape.element)
            # Insert into the new slide's shape tree
            target_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
            return
        
    raise ValueError(f"No table found on slide {source_slide_idx} to copy")