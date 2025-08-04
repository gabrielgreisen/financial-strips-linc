from pptx import Presentation
import pandas as pd
from pptx.dml.color import RGBColor
from helpers.copy_helpers import copy_table_from_template_slide
from helpers.logo_resources import get_logo_file_path
from helpers.logo_placement import place_logo_on_slide

def strips_layout_two(prs: Presentation, layout_index: int, buyers_chunk_df: pd.DataFrame, start_number: int):
    
    """
    Adds a slide to the presentation using the specified layout index,
    finds the table, and fills it with buyers from the DataFrame slice.

    Parameters
    ----------
    prs : pptx.Presentation
        The loaded Presentation object.

    layout_index : int
        Index of the slide master layout (e.g. 1).

    buyers_chunk_df : pd.DataFrame
        A slice of the DataFrame, typically up to 6 rows.
    """

    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    copy_table_from_template_slide(prs, source_slide_idx=2, target_slide=slide)

    # Find the table shape on the slide
    table = None
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            table_shape = shape
            break
    if table is None:
        raise ValueError(f'No table found on slide layout {layout_index}!')
    
    # Fill table rows
    for i, (_, row) in enumerate(buyers_chunk_df.iterrows()):
        row_idx = i + 1
        # Assumes that columns follow the correct format/order of the columns according to mask; fixed column positions
        exchange = str(row.iloc[3])
        ticker = str(row.iloc[4])
        country = str(row.iloc[5])
        description = str(row.iloc[6])
        br_presence = str(row.iloc[7])
        acquisition_count = str(row.iloc[8])
        acquisition_names = str(row.iloc[9])
        def format_number(val):
            try:
                return f"{float(val):,.2f}" # comma as thousand separator + 2 decimals
            except:
                return str(val)

        revenue = format_number(row.iloc[12])
        ebitda = format_number(row.iloc[13])
        market_cap = format_number(row.iloc[14])
        employees = format_number(row.iloc[15])

        # Build first column (numbering)
        number = start_number + i
        cell = table.cell(row_idx, 0)
        for para in cell.text_frame.paragraphs:
            if para.runs:
                para.runs[0].text = str(number)
                break
        else:
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = str(number)
    
        # Build second column (potential buyer)
        cell = table.cell(row_idx, 1)
        if exchange == "Private":
            found_index = None
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"({exchange})\n{country}" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"({exchange})\n{country}"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
        else:
            found_index = None    
            for idx, para in enumerate(cell.text_frame.paragraphs):
                if para.runs:
                    para.runs[0].text = f"({exchange}:{ticker})\n{country}" # Assumes same formating for the whole text
                    found_index = idx
                    break
            else:
                # if none of the paragraphs had a run text is added with default formating
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"({exchange}:{ticker})\n{country}"
                found_index = len(cell.text_frame.paragraphs) - 1

            for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
                if idx > found_index:
                    cell.text_frame._element.remove(extra_para._element)
                else:
                    break
            
            
        # Build third column (financials)
        cell = table.cell(row_idx, 2)
        found_index = None
        for idx, para in enumerate(cell.text_frame.paragraphs):
            if para.runs:
                para.runs[0].text = f"Revenue: {revenue}\nEBITDA: {ebitda}\nMarket Cap: {market_cap}\nTotal Debt: XXX\nFTE: {employees}" # Assumes same formating for the whole text
                found_index = idx
                break
        else:
            # if none of the paragraphs had a run text is added with default formating
            para = cell.text_frame.add_paragraph()
            run = para.add_run()
            run.text = f"Revenue: {revenue}\nEBITDA: {ebitda}\nMarket Cap: {market_cap}\nTotal Debt:XXX\nFTE: {employees}"
            found_index = len(cell.text_frame.paragraphs) - 1

        for idx, extra_para in reversed(list(enumerate(cell.text_frame.paragraphs))):
            if idx > found_index:
                cell.text_frame._element.remove(extra_para._element)
            else:
                break




        # Build fourth column (description)
        cell = table.cell(row_idx, 3)
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = description
        else:
            cell.text = description
        
        
        # Build fifth column (BR Presence)
        check_mark = "\u2713" # ✓
        cross_mark = "\u2718" # ✘

        if br_presence == "Yes":
            # Figure out how to put the check and the cross formating (its some sort of font)(solved for now)
            cell = table.cell(row_idx, 4)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.text = check_mark
                run.font.color.rgb = RGBColor(0,168,126)
            else:
                cell.text_frame.clear()
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = check_mark
                run.font.color.rgb = RGBColor(0,168,126)
        elif br_presence == "No":
            cell = table.cell(row_idx, 4)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.text = cross_mark
                run.font.color.rgb = RGBColor(192,0,0)
            else:
                cell.text_frame.clear()
                para = cell.text_frame.add_paragraph()
                run = para.add_run()
                run.text = cross_mark
                run.font.color.rgb = RGBColor(192,0,0)

        # Build sixth column (M&A History)
        cell = table.cell(row_idx, 5)
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = f"Acquisitions: {acquisition_count} \n \nCompanies: {acquisition_names}"
            cell.text_frame._element.remove(cell.text_frame.paragraphs[-1]._element)
        else:
            cell.text = f"Acquisitions: {acquisition_count} \n \nCompanies: {acquisition_names}"
            cell.text_frame._element.remove(cell.text_frame.paragraphs[-1]._element)
            

        # Add logos to the second column
        logo_file = get_logo_file_path(row)
        if logo_file:
            place_logo_on_slide(slide, table_shape, table, row_idx, 1, logo_file,
                                width_spacing=0.95, height_spacing=0.50, left_spacing=0.03, top_spacing=0.18)


